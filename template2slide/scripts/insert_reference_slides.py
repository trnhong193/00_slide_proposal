#!/usr/bin/env python3
"""
Insert reference slides from System_architecture.pptx and Available_Slide.pptx into generated presentation.

Usage:
    python insert_reference_slides.py <generated_pptx> <project_info_json> [output_pptx]

This script:
1. Inserts architecture template slide from System_architecture.pptx after the architecture diagram slide
2. Inserts Available slides (2-10) after slide 1
3. Inserts Available slides (11-25) after the last slide
4. Copies background from generated slides to inserted slides
"""

import sys
import json
import shutil
from pathlib import Path
from copy import deepcopy

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Error: python-pptx not installed. Install with: pip install python-pptx")
    sys.exit(1)


# Mapping deployment method to slide index in System_architecture.pptx
DEPLOYMENT_SLIDE_MAP = {
    'cloud': 1,  # Slide 2 (0-indexed: 1)
    'on-premise': 2,  # Slide 3 (0-indexed: 2)
    'on-prem': 2,  # Alias for on-premise
    'hybrid': 3,  # Slide 4 - Hybrid (AI Interference at site, Dashboard + Training Cloud)
    'hybrid-training-on-prem': 4,  # Slide 5 - Hybrid (AI Interference + Training at site, Dashboard Cloud)
    'hybrid-training-onprem': 4,  # Alias
    '4g-vpn-bridge': None,  # Not in template
    'vimov': None,  # Not in template
}


def get_deployment_method(project_info_path):
    """Extract deployment method from project_info JSON file"""
    try:
        with open(project_info_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # Handle both nested and flat structures
        if 'project_info' in data:
            project_info = data['project_info']
        else:
            project_info = data
        
        deployment_method = project_info.get('deployment_method', '').lower()
        # Normalize deployment method names
        if 'cloud' in deployment_method:
            return 'cloud'
        elif 'hybrid' in deployment_method:
            # Check if training is on-premise
            if 'training' in deployment_method and ('on-prem' in deployment_method or 'onprem' in deployment_method):
                return 'hybrid-training-on-prem'
            else:
                return 'hybrid'
        elif 'on-prem' in deployment_method or 'onprem' in deployment_method:
            return 'on-premise'
        elif '4g' in deployment_method or 'vpn' in deployment_method:
            return '4g-vpn-bridge'
        elif 'vimov' in deployment_method:
            return 'vimov'
        
        return deployment_method
    except Exception as e:
        print(f"Warning: Could not read deployment method from {project_info_path}: {e}")
        return None


def find_architecture_slide_index(presentation):
    """Find the index of the architecture diagram slide (type: diagram)"""
    # The architecture slide is typically the one with "System Architecture" or "Proposed System Architecture" title
    # We'll look for slides with diagram-related content
    for i, slide in enumerate(presentation.slides):
        # Check slide notes or title for architecture indicators
        # Since we can't easily check slide type, we'll use a heuristic:
        # Look for slides with "Architecture" in title or notes
        try:
            # Check if slide has a title shape
            for shape in slide.shapes:
                if hasattr(shape, 'text') and 'architecture' in shape.text.lower():
                    return i
        except:
            pass
    
    # Default: assume architecture slide is around slide 4 (0-indexed: 3)
    # This is based on SLIDE_TEMPLATE.md which says System Architecture is Slide 4
    return 3


def copy_slide_background(source_slide, target_slide):
    """Copy background from source slide to target slide"""
    try:
        # Get background from source slide
        source_bg = source_slide.background
        target_bg = target_slide.background
        
        # Copy background fill
        if source_bg.fill:
            target_bg.fill = deepcopy(source_bg.fill)
        
        # Copy background image if exists
        # Note: Background images are stored in slide relationships
        # This is a simplified approach - full implementation would need to handle relationships
        return True
    except Exception as e:
        print(f"Warning: Could not copy background: {e}")
        return False


def duplicate_slide(pres, source_index):
    """
    Duplicate a slide in the presentation (based on rearrange.py).
    Returns the new slide.
    """
    import six
    
    source = pres.slides[source_index]
    
    # Use source's layout to preserve formatting
    new_slide = pres.slides.add_slide(source.slide_layout)
    
    # Collect all image and media relationships from the source slide
    image_rels = {}
    for rel_id, rel in six.iteritems(source.part.rels):
        if "image" in rel.reltype or "media" in rel.reltype:
            image_rels[rel_id] = rel
    
    # CRITICAL: Clear placeholder shapes to avoid duplicates
    for shape in list(new_slide.shapes):
        try:
            sp = shape.element
            sp.getparent().remove(sp)
        except:
            pass
    
    # Copy all shapes from source
    for shape in source.shapes:
        try:
            el = shape.element
            new_el = deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
            
            # Handle picture shapes - need to update the blip reference
            blips = new_el.xpath(".//a:blip[@r:embed]")
            for blip in blips:
                old_rId = blip.get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                )
                if old_rId in image_rels:
                    # Create a new relationship in the destination slide for this image
                    old_rel = image_rels[old_rId]
                    new_rId = new_slide.part.rels.get_or_add(
                        old_rel.reltype, old_rel._target
                    )
                    # Update the blip's embed reference to use the new relationship ID
                    blip.set(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
                        new_rId,
                    )
        except Exception as e:
            print(f"Warning: Could not copy shape: {e}")
    
    # Copy background
    try:
        source_bg = source.background
        target_bg = new_slide.background
        
        if source_bg.fill:
            target_bg.fill = deepcopy(source_bg.fill)
    except Exception as e:
        print(f"Warning: Could not copy background: {e}")
    
    return new_slide


def copy_slide_from_other_pres(source_pres, source_index, target_pres):
    """
    Copy a slide from source presentation to target presentation.
    Returns the new slide.
    """
    import six
    
    source_slide = source_pres.slides[source_index]
    
    # Use source slide's layout
    new_slide = target_pres.slides.add_slide(source_slide.slide_layout)
    
    # Collect all image and media relationships from the source slide
    image_rels = {}
    for rel_id, rel in six.iteritems(source_slide.part.rels):
        if "image" in rel.reltype or "media" in rel.reltype:
            image_rels[rel_id] = rel
    
    # Clear existing shapes (placeholders)
    for shape in list(new_slide.shapes):
        try:
            sp = shape.element
            sp.getparent().remove(sp)
        except:
            pass
    
    # Copy all shapes from source slide
    for shape in source_slide.shapes:
        try:
            el = shape.element
            new_el = deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
            
            # Handle picture shapes - update blip references
            blips = new_el.xpath(".//a:blip[@r:embed]")
            for blip in blips:
                old_rId = blip.get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                )
                if old_rId in image_rels:
                    old_rel = image_rels[old_rId]
                    new_rId = new_slide.part.rels.get_or_add(old_rel.reltype, old_rel._target)
                    blip.set(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
                        new_rId,
                    )
        except Exception as e:
            print(f"Warning: Could not copy shape: {e}")
    
    # Copy background
    try:
        source_bg = source_slide.background
        target_bg = new_slide.background
        
        if source_bg.fill:
            target_bg.fill = deepcopy(source_bg.fill)
    except Exception as e:
        print(f"Warning: Could not copy background: {e}")
    
    return new_slide


def insert_reference_slides(generated_pptx_path, project_info_path, output_pptx_path=None):
    """
    Insert reference slides into generated presentation.
    
    Args:
        generated_pptx_path: Path to the generated PowerPoint file
        project_info_path: Path to project_info.json file
        output_pptx_path: Path for output file (default: overwrite input)
    """
    # Convert to absolute paths
    generated_pptx = Path(generated_pptx_path).resolve()
    project_info = Path(project_info_path).resolve()
    
    if not generated_pptx.exists():
        print(f"Error: Generated PPTX not found: {generated_pptx}")
        print(f"Current working directory: {Path.cwd()}")
        return False
    
    if not project_info.exists():
        print(f"Warning: Project info not found: {project_info}. Skipping architecture template slide.")
        deployment_method = None
    else:
        deployment_method = get_deployment_method(project_info)
    
    # Set output path
    if output_pptx_path is None:
        output_pptx_path = generated_pptx
    else:
        output_pptx_path = Path(output_pptx_path)
    
    # Paths to reference files
    script_dir = Path(__file__).parent
    ref_dir = script_dir.parent / "ref"  # ref directory is in template2slide folder
    system_arch_path = ref_dir / "System_architecture.pptx"
    # Try both possible file names
    available_slides_path = ref_dir / "AvailableSlide11.pptx"
    if not available_slides_path.exists():
        available_slides_path = ref_dir / "Available _Slide.pptx"
    
    if not system_arch_path.exists():
        print(f"Warning: System_architecture.pptx not found at {system_arch_path}")
        system_arch_path = None
    
    if not available_slides_path.exists():
        print(f"Warning: Available _Slide.pptx not found at {available_slides_path}")
        available_slides_path = None
    
    # Load presentations
    print(f"Loading generated presentation: {generated_pptx}")
    generated_pres = Presentation(str(generated_pptx))
    
    total_slides_before = len(generated_pres.slides)
    print(f"Original presentation has {total_slides_before} slides")
    
    # Create a copy for output
    if output_pptx_path != generated_pptx:
        shutil.copy2(generated_pptx, output_pptx_path)
        pres = Presentation(str(output_pptx_path))
    else:
        # Create a temporary file to work with
        temp_path = generated_pptx.parent / f"{generated_pptx.stem}_temp.pptx"
        shutil.copy2(generated_pptx, temp_path)
        pres = Presentation(str(temp_path))
    
    # Load reference presentations
    system_arch_pres = None
    available_pres = None
    
    if system_arch_path:
        system_arch_pres = Presentation(str(system_arch_path))
    if available_slides_path:
        available_pres = Presentation(str(available_slides_path))
    
    # Track slides to reorder
    slides_to_reorder = []
    
    # Step 1: Insert architecture template slide after architecture diagram slide
    arch_template_slide_index = None
    if system_arch_path and deployment_method:
        arch_slide_index = DEPLOYMENT_SLIDE_MAP.get(deployment_method)
        
        if arch_slide_index is not None and arch_slide_index < len(system_arch_pres.slides):
            print(f"\nStep 1: Inserting architecture template slide for deployment: {deployment_method}")
            print(f"  Using slide {arch_slide_index + 1} from System_architecture.pptx")
            
            # Find architecture diagram slide in generated presentation
            arch_diagram_index = find_architecture_slide_index(pres)
            print(f"  Architecture diagram slide found at index {arch_diagram_index}")
            
            # Copy slide from system_arch_pres
            new_slide = copy_slide_from_other_pres(system_arch_pres, arch_slide_index, pres)
            arch_template_slide_index = len(pres.slides) - 1
            slides_to_reorder.append((arch_template_slide_index, arch_diagram_index + 1))
            print(f"  ✓ Architecture template slide copied (will be moved to position {arch_diagram_index + 2})")
        else:
            print(f"\nStep 1: Skipping architecture template slide (deployment method '{deployment_method}' not in template)")
    else:
        print("\nStep 1: Skipping architecture template slide (file or deployment method not found)")
    
    # Step 2: Insert Available slides 2-10 after slide 1 (continuously, in order)
    available_slides_2_10_indices = []
    if available_pres and len(available_pres.slides) >= 10:
        print(f"\nStep 2: Inserting Available slides 2-10 after slide 1 (continuously)")
        # Insert all slides first, then reorder them to positions 2-10 in order
        for i in range(1, 10):  # Slides 2-10 from Available (indices 1-9)
            new_slide = copy_slide_from_other_pres(available_pres, i, pres)
            new_index = len(pres.slides) - 1
            # Target position: 2, 3, 4, ..., 10 (after slide 1)
            target_position = 1 + i  # 2, 3, 4, ..., 10
            available_slides_2_10_indices.append((new_index, target_position))
            print(f"  ✓ Copied Available slide {i + 1} (will be moved to position {target_position})")
        print(f"  ✓ All Available slides 2-10 copied (will be inserted continuously after slide 1)")
    else:
        print("\nStep 2: Skipping Available slides 2-10 (file not found or not enough slides)")
    
    # Step 3: Insert Available slides 11-25 at the end and copy background from generated slides
    if available_pres and len(available_pres.slides) >= 25:
        print(f"\nStep 3: Inserting Available slides 11-25 after last slide")
        # Get background from first generated slide (slide 0) to apply to inserted slides
        source_bg_slide = pres.slides[0] if len(pres.slides) > 0 else None
        
        for i in range(10, 25):  # Slides 11-25 (indices 10-24)
            new_slide = copy_slide_from_other_pres(available_pres, i, pres)
            # Copy background from generated slides to inserted slide
            if source_bg_slide:
                try:
                    source_bg = source_bg_slide.background
                    target_bg = new_slide.background
                    if source_bg.fill:
                        target_bg.fill = deepcopy(source_bg.fill)
                    print(f"  ✓ Copied Available slide {i + 1} (with background from generated slides)")
                except Exception as e:
                    print(f"  ✓ Copied Available slide {i + 1} (background copy failed: {e})")
            else:
                print(f"  ✓ Copied Available slide {i + 1}")
        print(f"  ✓ All Available slides 11-25 copied")
    else:
        print("\nStep 3: Skipping Available slides 11-25 (file not found or not enough slides)")
    
    # Reorder slides: move slides to correct positions
    print(f"\nReordering slides...")
    slides = pres.slides._sldIdLst
    
    # Collect all reorder operations
    all_reorders = []
    if arch_template_slide_index is not None and len(slides_to_reorder) > 0:
        all_reorders.append((arch_template_slide_index, slides_to_reorder[0][1]))
    
    # For Available slides 2-10: sort by target position (ascending) to insert in order
    # This ensures they are inserted continuously: 2, 3, 4, ..., 10
    available_slides_sorted = sorted(available_slides_2_10_indices, key=lambda x: x[1])
    all_reorders.extend(available_slides_sorted)
    
    # Sort by source index (descending) so we move from end to beginning
    # This prevents index shifting issues when moving multiple slides
    all_reorders.sort(key=lambda x: x[0], reverse=True)
    
    # Get current slide elements
    slide_elements = list(slides)
    
    for source_idx, target_idx in all_reorders:
        if source_idx < len(slide_elements):
            slide_element = slide_elements[source_idx]
            # Remove from current position
            slides.remove(slide_element)
            # Insert at target position (clamp to valid range)
            insert_pos = min(target_idx, len(slides))
            slides.insert(insert_pos, slide_element)
            print(f"  ✓ Moved slide from position {source_idx + 1} to position {target_idx + 1}")
    
    # Save final presentation
    if output_pptx_path == generated_pptx:
        pres.save(str(generated_pptx))
        if temp_path.exists():
            temp_path.unlink()
    else:
        pres.save(str(output_pptx_path))
    
    total_slides_after = len(pres.slides)
    print(f"\n✓ Complete! Presentation now has {total_slides_after} slides (was {total_slides_before})")
    print(f"  Saved to: {output_pptx_path}")
    
    return True


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python insert_reference_slides.py <generated_pptx> <project_info_json> [output_pptx]")
        print("\nExample:")
        print("  python insert_reference_slides.py output/proposal.pptx output/project_info.json output/final_proposal.pptx")
        sys.exit(1)
    
    generated_pptx = sys.argv[1]
    project_info = sys.argv[2]
    output_pptx = sys.argv[3] if len(sys.argv) > 3 else None
    
    success = insert_reference_slides(generated_pptx, project_info, output_pptx)
    sys.exit(0 if success else 1)

